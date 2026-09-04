#!/usr/bin/env python3
"""Build the editable eight-slide Bangladesh Innovation Fair pitch deck."""

from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "output" / "presentation" / "digital-delta-innovation-fair.pptx"

W = 13.333
H = 7.5
INK = "0A2638"
DARK = "061A27"
TEAL = "087F8C"
CYAN = "1EB8BC"
CORAL = "EF5B5B"
AMBER = "E9A928"
GREEN = "2E9D72"
LIGHT = "F4F8F9"
PALE = "E5F1F2"
PALE_BLUE = "EAF1F6"
MID = "597181"
LINE = "CCD9DF"
WHITE = "FFFFFF"
FONT = "Aptos"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def rect(slide, x, y, w, h, fill, radius=True, line="none", transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    if line == "none":
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    return shape


def text_box(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=18,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0,
    italic=False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def rich_text(slide, runs, x, y, w, h, size=18, color=INK, margin=0, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.space_after = Pt(0)
    for item in runs:
        run = p.add_run()
        run.text = item[0]
        run.font.name = FONT
        run.font.size = Pt(item[1] if len(item) > 1 else size)
        run.font.bold = item[2] if len(item) > 2 else False
        run.font.color.rgb = rgb(item[3] if len(item) > 3 else color)
    return box


def add_title(slide, number, title, subtitle=None, dark=False):
    main = WHITE if dark else INK
    accent = CYAN if dark else TEAL
    text_box(slide, f"0{number}", 0.58, 0.34, 0.52, 0.34, 12, accent, True)
    text_box(slide, title, 1.12, 0.23, 11.35, 0.55, 27, main, True)
    rect(slide, 0.58, 0.92, 0.68, 0.055, accent, radius=False)
    if subtitle:
        text_box(slide, subtitle, 1.42, 0.82, 10.7, 0.38, 12, MID if not dark else "A9C1CC")


def add_footer(slide, number, dark=False, label="DIGITAL DELTA"):
    color = "8CA4AF" if dark else MID
    text_box(slide, label, 0.58, 7.13, 2.2, 0.2, 8, color, True)
    text_box(slide, str(number), 12.26, 7.11, 0.45, 0.22, 9, color, True, PP_ALIGN.RIGHT)


def add_notes(slide, notes: str):
    frame = slide.notes_slide.notes_text_frame
    frame.text = notes


def set_picture_alt(shape, description: str):
    c_nv_pr = shape._element.xpath(".//p:cNvPr")[0]
    c_nv_pr.set("descr", description)
    c_nv_pr.set("title", description[:120])


def add_picture_cover(slide, path: Path, x, y, w, h, alt: str):
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    frame_ratio = w / h
    if image_ratio > frame_ratio:
        shown_width = h * image_ratio
        crop = (shown_width - w) / shown_width / 2
        shape = slide.shapes.add_picture(str(path), Inches(x), Inches(y), height=Inches(h))
        shape.crop_left = crop
        shape.crop_right = crop
    else:
        shown_height = w / image_ratio
        crop = (shown_height - h) / shown_height / 2
        shape = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
        shape.crop_top = crop
        shape.crop_bottom = crop
    shape.left, shape.top, shape.width, shape.height = Inches(x), Inches(y), Inches(w), Inches(h)
    set_picture_alt(shape, alt)
    return shape


def add_picture_fit(slide, path: Path, x, y, w, h, alt: str):
    with Image.open(path) as image:
        image_ratio = image.width / image.height
    frame_ratio = w / h
    if image_ratio > frame_ratio:
        width = w
        height = w / image_ratio
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * image_ratio
        top = y
        left = x + (w - width) / 2
    shape = slide.shapes.add_picture(str(path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
    set_picture_alt(shape, alt)
    return shape


def bullet_list(slide, items, x, y, w, h, size=17, color=INK, accent=TEAL, gap=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0)
    frame.margin_right = Inches(0)
    frame.margin_top = Inches(0)
    frame.margin_bottom = Inches(0)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.space_after = Pt(gap)
        p.level = 0
        p.text = f"•  {item}"
        p.runs[0].font.name = FONT
        p.runs[0].font.size = Pt(size)
        p.runs[0].font.color.rgb = rgb(color)
    return box


def stat(slide, value, label, x, y, w, fill=PALE, value_color=TEAL):
    rect(slide, x, y, w, 1.02, fill, line=LINE)
    text_box(slide, value, x + 0.16, y + 0.11, w - 0.32, 0.43, 24, value_color, True)
    text_box(slide, label, x + 0.16, y + 0.59, w - 0.32, 0.25, 10.5, MID, True)


def arrow(slide, x1, y1, x2, y2, color=TEAL):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(2.5)
    line.line.end_arrowhead = True
    return line


def new_slide(prs, fill=WHITE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(fill)
    return slide


def build_cover(prs):
    slide = new_slide(prs, DARK)
    add_picture_cover(
        slide,
        ROOT / "artifacts/screenshots/command-en/02-command-en-live-overview-1920x1080.png",
        6.8,
        0,
        6.533,
        7.5,
        "English headquarters overview with the live Sylhet map and operations panels",
    )
    overlay = rect(slide, 6.3, 0, 1.6, 7.5, DARK, radius=False, transparency=18)
    overlay.rotation = 0
    rect(slide, 0.62, 0.64, 1.16, 0.08, CYAN, radius=False)
    text_box(slide, "BANGLADESH INNOVATION FAIR 2026", 0.62, 0.85, 5.45, 0.35, 12, CYAN, True)
    text_box(slide, "Digital\nDelta", 0.62, 1.42, 5.65, 1.45, 45, WHITE, True)
    text_box(
        slide,
        "Relief logistics that keep working when the internet and power do not",
        0.62,
        3.06,
        5.55,
        1.05,
        23,
        "D5E7ED",
        False,
    )
    rect(slide, 0.62, 4.55, 5.55, 1.26, "0C3448", line="244E61")
    rich_text(
        slide,
        [
            ("FIELD FIRST  ", 11, True, CYAN),
            ("Bangla + English  •  Android  •  Protobuf  •  Signed custody", 14, True, WHITE),
        ],
        0.9,
        4.87,
        5.0,
        0.58,
    )
    text_box(slide, "Touhidul Alam Seyam  •  Individual project lead and developer", 0.62, 6.64, 5.5, 0.28, 10.5, "9FB7C2")
    add_footer(slide, 1, dark=True)
    add_notes(slide, "Opening: This is not another cloud dashboard. The field mission lives on the phones. The projector only observes it.")


def build_problem(prs):
    slide = new_slide(prs)
    add_title(slide, 2, "One flood can break every dependency at once", "The continuity gap is local, not abstract")
    stat(slide, "7.2M", "people affected in the 2022 northeast floods", 0.65, 1.35, 2.35, PALE_BLUE, CORAL)
    stat(slide, "3.74M", "people affected across six districts in 2024", 3.16, 1.35, 2.35, PALE_BLUE, CORAL)
    stat(slide, "1.4M", "needed urgent assistance in the 2024 plan", 5.67, 1.35, 2.35, PALE_BLUE, CORAL)
    rect(slide, 0.65, 2.72, 7.37, 3.56, LIGHT, line=LINE)
    text_box(slide, "When the flood arrives", 0.95, 2.99, 3.05, 0.4, 18, INK, True)
    bullet_list(
        slide,
        [
            "Electricity may be shut down for safety.",
            "Mobile communication becomes difficult.",
            "Roads, storage, and airport access can fail together.",
            "The request, route, and handoff record still have to move.",
        ],
        0.95,
        3.53,
        6.55,
        2.35,
        16,
        INK,
        gap=13,
    )
    rect(slide, 8.35, 1.35, 4.3, 4.93, DARK, line=DARK)
    text_box(slide, "THE GAP", 8.72, 1.73, 3.56, 0.3, 11, CYAN, True)
    text_box(slide, "Central systems can be excellent and still be unreachable.", 8.72, 2.14, 3.5, 1.18, 24, WHITE, True)
    text_box(slide, "Field teams need continuity beside existing radio, phone, and paper procedures, not a replacement for them.", 8.72, 3.66, 3.45, 1.43, 16, "C6D9E1")
    text_box(slide, "Problem evidence from DDM and UN Bangladesh", 8.72, 5.57, 3.35, 0.34, 10, "86A9B8", True)
    add_footer(slide, 2)
    add_notes(slide, "Sources: https://rapid.ddm.gov.bd/risk/riskinfo ; https://bangladesh.un.org/en/188010-flash-floods-humanitarian-response-plan-2022-united-nations-bangladesh-coordinated-appeal ; https://bangladesh.un.org/sites/default/files/2024-10/Bangladesh-HCTT-Humanitarian%20Response%20Plan-Cyclone%20and%20Monsoon%20Floods%202024-29-09-2024%20%28June%202024%20to%20March%202025%29.pdf")


def build_architecture(prs):
    slide = new_slide(prs)
    add_title(slide, 3, "The mission lives on the phones", "The observer can disappear without stopping field work")
    nodes = [
        (0.65, "CLINIC", "Create P0 request", CORAL),
        (3.22, "RELAY", "Store before ack", TEAL),
        (5.79, "HOSPITAL", "Decrypt + decide", GREEN),
    ]
    for x, label, body, color in nodes:
        rect(slide, x, 1.45, 2.05, 1.33, WHITE, line=color)
        text_box(slide, label, x + 0.16, 1.67, 1.72, 0.28, 12, color, True, PP_ALIGN.CENTER)
        text_box(slide, body, x + 0.16, 2.08, 1.72, 0.33, 13, INK, True, PP_ALIGN.CENTER)
    arrow(slide, 2.73, 2.12, 3.12, 2.12)
    arrow(slide, 5.3, 2.12, 5.69, 2.12)
    text_box(slide, "AES-256-GCM payload  •  RSA-OAEP key wrap  •  Protobuf envelope", 1.22, 3.05, 6.08, 0.34, 13, TEAL, True, PP_ALIGN.CENTER)
    rect(slide, 8.45, 1.45, 4.15, 3.72, PALE_BLUE, line=LINE)
    text_box(slide, "OPTIONAL OBSERVER PATH", 8.78, 1.72, 3.48, 0.28, 11, MID, True)
    text_box(slide, "Go observer", 8.78, 2.14, 1.25, 0.35, 16, INK, True)
    text_box(slide, "ordered, sanitized events", 10.18, 2.17, 2.02, 0.3, 12, MID)
    arrow(slide, 9.34, 2.78, 9.34, 3.25, TEAL)
    text_box(slide, "Next.js headquarters", 8.78, 3.39, 2.35, 0.35, 16, INK, True)
    text_box(slide, "local SSE + offline map", 8.78, 3.83, 2.92, 0.3, 12, MID)
    text_box(slide, "Vercel + D1", 8.78, 4.45, 1.55, 0.3, 13, INK, True)
    text_box(slide, "reviewer convenience only", 10.2, 4.46, 1.98, 0.3, 11.5, MID)
    rect(slide, 0.65, 4.02, 7.37, 1.15, PALE, line=TEAL)
    rich_text(
        slide,
        [
            ("DISCONNECT THE LAPTOP  ", 13, True, TEAL),
            ("Requests, relay, routing, triage, and custody continue on Android.", 17, True, INK),
        ],
        0.94,
        4.35,
        6.78,
        0.5,
    )
    text_box(slide, "Field transport", 0.65, 5.78, 1.45, 0.28, 10, MID, True)
    text_box(slide, "Nearby framed Protobuf", 2.15, 5.75, 2.3, 0.3, 14, INK, True)
    text_box(slide, "Supported IP link", 5.0, 5.78, 1.3, 0.28, 10, MID, True)
    text_box(slide, "gRPC + Protobuf", 6.38, 5.75, 1.7, 0.3, 14, INK, True)
    text_box(slide, "Browser boundary", 8.45, 5.78, 1.35, 0.28, 10, MID, True)
    text_box(slide, "Sanitized SSE JSON", 9.92, 5.75, 2.35, 0.3, 14, INK, True)
    add_footer(slide, 3)
    add_notes(slide, "Architecture evidence: README.md and docs/ARCHITECTURE.md. JSON is used only at the browser observer boundary, never on the mesh.")


def build_modules(prs):
    slide = new_slide(prs)
    add_title(slide, 4, "Eight modules, one continuous mission", "Each module changes the same signed and replayable field state")
    modules = [
        ("M1", "Identity", "Offline role credential + audit", TEAL),
        ("M2", "CRDT sync", "Vector clocks + human conflict", CYAN),
        ("M3", "Mesh relay", "Authenticated store and forward", AMBER),
        ("M4", "Routing", "Road + water + simulated air", GREEN),
        ("M5", "Custody", "Signed QR + replay rejection", CORAL),
        ("M6", "Triage", "SLA warning + safe preemption", TEAL),
        ("M7", "Route risk", "On-device ONNX advisory", AMBER),
        ("M8", "Hybrid fleet", "Rendezvous + signed handoff", CYAN),
    ]
    for index, (code, label, body, color) in enumerate(modules):
        col = index % 4
        row = index // 4
        x = 0.65 + col * 3.13
        y = 1.42 + row * 2.34
        rect(slide, x, y, 2.84, 1.93, WHITE, line=LINE)
        rect(slide, x, y, 0.12, 1.93, color, radius=False)
        text_box(slide, code, x + 0.28, y + 0.24, 0.55, 0.3, 12, color, True)
        text_box(slide, label, x + 0.84, y + 0.19, 1.68, 0.4, 18, INK, True)
        text_box(slide, body, x + 0.28, y + 0.83, 2.28, 0.62, 14, MID)
    rect(slide, 0.65, 6.25, 12.0, 0.56, DARK, line=DARK)
    text_box(slide, "REAL LOGIC", 0.92, 6.42, 1.05, 0.2, 10, CYAN, True)
    text_box(slide, "Simulated weather and vehicles are labelled. Cryptography, storage, routing, and decisions are real code paths.", 2.05, 6.35, 10.05, 0.3, 13.5, WHITE, True)
    add_footer(slide, 4)
    add_notes(slide, "Technical capability: all eight module paths are implemented. Simulation labels remain visible for environmental inputs and vehicle movement.")


def build_mission(prs):
    slide = new_slide(prs)
    add_title(slide, 5, "One failure changes route and priority", "A deterministic offline drill that judges can follow in seconds")
    phone_paths = [
        (ROOT / "artifacts/screenshots/field-en/14-field-en-route-rerouted-1280x2856.png", "Offline route screen showing simulated E3 failure and boat reroute"),
        (ROOT / "artifacts/screenshots/field-en/15-field-en-triage-preemption-1280x2856.png", "Triage screen showing predicted P0 SLA breach and proposed preemption"),
        (ROOT / "artifacts/screenshots/field-en/16-field-en-triage-confirmed-1280x2856.png", "Triage screen showing coordinator-confirmed P2 deposit"),
    ]
    for i, (path, alt) in enumerate(phone_paths):
        x = 0.5 + i * 2.22
        rect(slide, x, 1.33, 2.02, 4.94, WHITE, line=LINE)
        add_picture_fit(slide, path, x + 0.08, 1.41, 1.86, 4.76, alt)
    steps = [
        ("01", "E3 fails", "Truck path becomes invalid", CORAL),
        ("02", "Boat selected", "E6 + E7, ETA 200 min", TEAL),
        ("03", "P0 warning", "30% slowdown breaches SLA", AMBER),
        ("04", "Human confirms", "P2 deposit at N3", GREEN),
    ]
    for index, (num, title, body, color) in enumerate(steps):
        y = 1.38 + index * 1.25
        rect(slide, 7.55, y, 4.93, 0.98, LIGHT, line=LINE)
        text_box(slide, num, 7.78, y + 0.17, 0.44, 0.28, 12, color, True)
        text_box(slide, title, 8.35, y + 0.12, 1.65, 0.34, 16, INK, True)
        text_box(slide, body, 10.0, y + 0.17, 2.1, 0.3, 12.5, MID)
    rect(slide, 7.55, 6.44, 4.93, 0.42, PALE, line=TEAL)
    text_box(slide, "SIMULATED FAILURE  •  REAL ROUTE + TRIAGE", 7.75, 6.55, 4.53, 0.18, 9.5, TEAL, True, PP_ALIGN.CENTER)
    add_footer(slide, 5)
    add_notes(slide, "Demo sequence: show the seeded truck route, inject simulated E3 failure, show boat E6 and E7, show P0 slowdown breach, then require coordinator confirmation before changing the P2 assignment.")


def build_security(prs):
    slide = new_slide(prs)
    add_title(slide, 6, "Trust is verified at every handoff", "The relay can carry the envelope but cannot read its contents")
    add_picture_fit(
        slide,
        ROOT / "artifacts/screenshots/field-en/17-field-en-pod-verified.png",
        0.5,
        1.33,
        2.35,
        4.92,
        "Proof-of-delivery screen showing the first signed handoff accepted",
    )
    add_picture_fit(
        slide,
        ROOT / "artifacts/screenshots/field-en/18-field-en-pod-replay-rejected.png",
        2.92,
        1.33,
        2.35,
        4.92,
        "Proof-of-delivery screen showing replayed nonce rejection",
    )
    text_box(slide, "Security path", 5.72, 1.38, 2.2, 0.34, 18, INK, True)
    bullet_list(
        slide,
        [
            "Administrator-signed role credentials",
            "Fresh peer challenge + RSA-PSS proof",
            "RSA-OAEP wrapped AES-256-GCM payload",
            "Atomic nonce claim + linked receipt hash",
            "Signed durable acknowledgements",
        ],
        5.72,
        1.9,
        3.3,
        3.45,
        14.5,
        INK,
        gap=10,
    )
    stat(slide, "60", "connected Android journeys in the current suite", 9.35, 1.42, 2.95, PALE_BLUE, TEAL)
    stat(slide, "10,000", "gRPC streams durably acknowledged in the load run", 9.35, 2.7, 2.95, PALE_BLUE, TEAL)
    stat(slide, "17", "headquarters tests plus production build", 9.35, 3.98, 2.95, PALE_BLUE, TEAL)
    rect(slide, 5.72, 5.67, 6.58, 0.72, "FFF3E3", line=AMBER)
    text_box(slide, "HONEST LIMIT", 5.95, 5.87, 1.12, 0.22, 10, AMBER, True)
    text_box(slide, "Physical three-phone and camera-in-airplane-mode proof remains a release gate.", 7.2, 5.79, 4.82, 0.38, 13, INK, True)
    add_footer(slide, 6)
    add_notes(slide, "Evidence: docs/TESTING.md and artifacts/reports/load/2026-09-04-go-10000.md. The 10,000-stream run proves connection capacity under recorded laptop conditions, not district radio throughput. QR camera and three-phone evidence remain open.")


def build_market(prs):
    slide = new_slide(prs)
    add_title(slide, 7, "Start with one district exercise", "Public-interest deployment, not a consumer subscription")
    table = slide.shapes.add_table(5, 3, Inches(0.65), Inches(1.4), Inches(7.05), Inches(3.5)).table
    table.columns[0].width = Inches(1.75)
    table.columns[1].width = Inches(2.18)
    table.columns[2].width = Inches(3.12)
    rows = [
        ("Alternative", "Strong at", "Remaining continuity gap"),
        ("Sahana Eden", "Relief logistics", "Server reachability in the field"),
        ("KoboCollect", "Offline forms", "Routing, relay, signed custody"),
        ("Briar", "Nearby secure messaging", "Operational logistics workflow"),
        ("Meshtastic", "Off-grid radio mesh", "Requires additional hardware"),
    ]
    for r_index, row in enumerate(rows):
        for c_index, value in enumerate(row):
            cell = table.cell(r_index, c_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(INK if r_index == 0 else (PALE_BLUE if r_index % 2 else WHITE))
            cell.margin_left = Inches(0.11)
            cell.margin_right = Inches(0.11)
            cell.margin_top = Inches(0.08)
            cell.margin_bottom = Inches(0.07)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.text = value
            p.runs[0].font.name = FONT
            p.runs[0].font.size = Pt(11.5 if r_index else 11)
            p.runs[0].font.bold = r_index == 0 or c_index == 0
            p.runs[0].font.color.rgb = rgb(WHITE if r_index == 0 else INK)
    rect(slide, 8.05, 1.4, 4.58, 3.5, LIGHT, line=LINE)
    text_box(slide, "Beachhead", 8.37, 1.7, 1.6, 0.33, 18, TEAL, True)
    text_box(slide, "Six northeastern\ndistricts", 8.37, 2.2, 3.55, 0.86, 21, INK, True)
    text_box(slide, "about 64 upazilas  •  360 unions", 8.37, 3.08, 3.55, 0.3, 13, MID)
    text_box(slide, "First obtainable deployment", 8.37, 3.61, 2.8, 0.28, 11, MID, True)
    text_box(slide, "One district  •  one health referral path  •  one relief partner", 8.37, 4.0, 3.74, 0.62, 16, INK, True)
    labels = [
        ("1", "Map the workflow", "Authority + data owner"),
        ("2", "Controlled exercise", "Range + recovery + usability"),
        ("3", "Shadow pilot", "Compare with official log"),
    ]
    for index, (num, title, body) in enumerate(labels):
        x = 0.65 + index * 4.0
        rect(slide, x, 5.28, 3.72, 1.2, WHITE, line=TEAL)
        text_box(slide, num, x + 0.2, 5.56, 0.35, 0.32, 18, TEAL, True)
        text_box(slide, title, x + 0.67, 5.47, 2.35, 0.56, 14.5, INK, True)
        text_box(slide, body, x + 0.67, 6.08, 2.75, 0.25, 10.8, MID)
    add_footer(slide, 7)
    add_notes(slide, "Sources: https://sahanafoundation.org/products/eden/ ; https://support.kobotoolbox.org/data_collection_kobocollect.html ; https://briarproject.org/manual/ ; https://meshtastic.org/docs/ ; Bangladesh 2024 to 2025 humanitarian response plan. Commercialization is deployment preparation, integration, training, offline map packaging, exercise support, and maintained releases.")


def build_close(prs):
    slide = new_slide(prs, DARK)
    add_title(slide, 8, "The next milestone is a real field exercise", "Software is ready to be challenged by phones, radios, people, and procedure", dark=True)
    rect(slide, 0.65, 1.48, 6.6, 1.42, "0C3448", line="28576A")
    text_box(slide, "THE ASK", 0.94, 1.75, 0.92, 0.25, 11, CYAN, True)
    text_box(slide, "Run one controlled three-phone exercise with a district authority, health facility, and relief partner", 1.98, 1.65, 4.88, 0.98, 17, WHITE, True)
    text_box(slide, "Success gates", 0.65, 3.2, 2.45, 0.4, 18, WHITE, True)
    bullet_list(
        slide,
        [
            "A to B to C relay with the middle phone restarted",
            "Camera provisioning and custody in airplane mode",
            "Target-phone memory, route latency, and battery",
            "Bangla, English, TalkBack, and operator task review",
        ],
        0.65,
        3.75,
        6.08,
        2.4,
        15,
        "D8E7EC",
        gap=10,
    )
    rect(slide, 7.62, 1.48, 4.98, 4.82, WHITE, line=WHITE)
    add_picture_fit(
        slide,
        ROOT / "artifacts/screenshots/field-en/24-field-en-hybrid-replanned-1280x2856.png",
        7.86,
        1.68,
        2.0,
        4.37,
        "English hybrid-fleet screen showing delayed boat replanning from R3 to R2",
    )
    text_box(slide, "Built solo by", 10.18, 2.02, 1.65, 0.25, 11, MID, True)
    text_box(slide, "Touhidul\nAlam Seyam", 10.18, 2.38, 2.02, 0.92, 21, INK, True)
    text_box(slide, "Product  •  Android  •  Go\nSecurity  •  Routing  •  ML\nNext.js  •  Testing  •  Evidence", 10.18, 3.58, 2.0, 1.13, 13, MID)
    text_box(slide, "SDG 3  •  9  •  11  •  13  •  16", 10.18, 5.21, 1.98, 0.28, 10.5, TEAL, True)
    text_box(slide, "digital-delta-headquarters.vercel.app", 0.65, 6.61, 4.25, 0.26, 11, CYAN, True)
    text_box(slide, "github.com/Seyamalam/digital-delta", 5.17, 6.61, 4.0, 0.26, 11, "9FB7C2", True)
    text_box(slide, "Thank you", 10.15, 6.53, 2.1, 0.4, 20, WHITE, True, PP_ALIGN.RIGHT)
    add_footer(slide, 8, dark=True)
    add_notes(slide, "Close: The fair build proves the hard software boundary. The ask is not national deployment. It is one controlled exercise that can turn emulator evidence into measured field evidence. Inclusion and ethics gates are explicit in docs/PHYSICAL_DEVICE_TEST.md and docs/FINAL_REPORT.md.")


def main() -> int:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "Digital Delta innovation fair pitch"
    prs.core_properties.subject = "Offline disaster logistics and verified relief delivery in Bangladesh"
    prs.core_properties.author = "Touhidul Alam Seyam"
    prs.core_properties.last_modified_by = "Touhidul Alam Seyam"
    prs.core_properties.comments = "Editable eight-slide fair submission deck"
    build_cover(prs)
    build_problem(prs)
    build_architecture(prs)
    build_modules(prs)
    build_mission(prs)
    build_security(prs)
    build_market(prs)
    build_close(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)
    return 0


if __name__ == "__main__":
    sys.exit(main())
